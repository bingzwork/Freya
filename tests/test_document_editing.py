from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.capabilities.registration_bridge import CapabilityRegistrationBridge
from app.capabilities.router import CapabilityRouter
from app.core.file_allowlist import AccessRule, FileAllowlist, FileAllowlistConfig, FileOperation
from app.core.tool_manager import ToolManager
from app.document.capability import DocumentEditingCapability
from app.orchestrator.capabilities import create_all_capabilities
from app.orchestrator.capability_registry import CapabilityRegistry, reset_capability_registry
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


def allowlist(root: Path) -> FileAllowlist:
    config = FileAllowlistConfig(allow_current_directory=False, allow_temp_directory=False, allow_home_directory=False)
    config.allowed_extensions.update({".txt", ".md", ".markdown", ".html", ".htm", ".docx", ".pdf", ".xlsx", ".csv", ".pptx"})
    result = FileAllowlist(config)
    result.add_rule(AccessRule(pattern=f"{root.resolve()}/**", operations={FileOperation.READ, FileOperation.WRITE, FileOperation.CREATE}, description="document test workspace"))
    return result


def capability(tmp_path: Path) -> DocumentEditingCapability:
    return DocumentEditingCapability(file_allowlist=allowlist(tmp_path), output_root=tmp_path)


def test_markdown_edit_save_and_original_preservation(tmp_path: Path):
    source = tmp_path / "guide.md"
    source.write_text("# Intro\nOld introduction", encoding="utf-8")
    original = source.read_bytes()
    result = capability(tmp_path).execute("edit", {"file_path": str(source), "edits": [{"old": "Old introduction", "new": "Rewritten introduction"}]})
    assert result["success"] is True
    output = Path(result["saved_path"])
    assert output.read_text(encoding="utf-8") == "# Intro\nRewritten introduction"
    assert source.read_bytes() == original


def test_docx_edit_and_save(tmp_path: Path):
    source = tmp_path / "input.docx"
    document = Document(); document.add_paragraph("Original introduction"); document.save(source)
    result = capability(tmp_path).execute("edit", {"file_path": str(source), "edits": [{"old": "Original introduction", "new": "Rewritten introduction"}]})
    assert result["success"] is True
    edited = Document(result["saved_path"])
    assert edited.paragraphs[0].text == "Rewritten introduction"
    assert source.exists()


def test_spreadsheet_cell_edit_preserves_formula_and_sheets(tmp_path: Path):
    source = tmp_path / "numbers.xlsx"
    workbook = Workbook(); sheet = workbook.active; sheet.title = "Data"; sheet["A1"] = 2; sheet["A2"] = 3; sheet["A3"] = "=SUM(A1:A2)"; workbook.create_sheet("Keep") ; workbook.save(source)
    result = capability(tmp_path).execute("edit", {"file_path": str(source), "cells": [{"sheet": "Data", "cell": "A1", "value": 5}]})
    assert result["success"] is True
    edited = load_workbook(result["saved_path"], data_only=False)
    assert edited.sheetnames == ["Data", "Keep"]
    assert edited["Data"]["A1"].value == 5
    assert edited["Data"]["A3"].value == "=SUM(A1:A2)"


def test_pptx_text_edit_preserves_slide_count(tmp_path: Path):
    source = tmp_path / "deck.pptx"
    presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[6]); box = slide.shapes.add_textbox(100, 100, 500, 100); box.text = "Old title"; presentation.slides.add_slide(presentation.slide_layouts[6]); presentation.save(source)
    result = capability(tmp_path).execute("edit", {"file_path": str(source), "edits": [{"old": "Old title", "new": "New title"}]})
    assert result["success"] is True
    edited = Presentation(result["saved_path"])
    assert len(edited.slides) == 2
    assert edited.slides[0].shapes[0].text == "New title"


def test_pdf_remove_page_and_export_text(tmp_path: Path):
    source = tmp_path / "input.pdf"
    output = source.open("wb"); pdf = canvas.Canvas(output, pagesize=letter); pdf.drawString(50, 750, "Page one"); pdf.showPage(); pdf.drawString(50, 750, "Page two"); pdf.save(); output.close()
    result = capability(tmp_path).execute("edit", {"file_path": str(source), "remove_pages": [1]})
    assert result["success"] is True
    assert len(PdfReader(result["saved_path"]).pages) == 1
    text_result = capability(tmp_path).execute("export", {"file_path": str(source), "target_extension": ".txt"})
    assert text_result["success"] is True
    assert "Page one" in Path(text_result["saved_path"]).read_text(encoding="utf-8")


def test_html_export_and_inspect(tmp_path: Path):
    source = tmp_path / "notes.txt"; source.write_text("<hello>", encoding="utf-8")
    result = capability(tmp_path).execute("export", {"file_path": str(source), "target_extension": ".html"})
    assert result["success"] is True
    assert "&lt;hello&gt;" in Path(result["saved_path"]).read_text(encoding="utf-8")
    inspection = capability(tmp_path).execute("inspect", {"file_path": str(source)})
    assert inspection["success"] is True
    assert inspection["inspection"]["character_count"] == 7


def test_unsupported_and_corrupt_files_fail_without_output(tmp_path: Path):
    unsupported = tmp_path / "input.exe"; unsupported.write_bytes(b"bad")
    result = capability(tmp_path).execute("inspect", {"file_path": str(unsupported)})
    assert result["success"] is False
    corrupt = tmp_path / "broken.docx"; corrupt.write_bytes(b"not a zip")
    result = capability(tmp_path).execute("edit", {"file_path": str(corrupt), "edits": []})
    assert result["success"] is False
    assert not (tmp_path / "broken.edited.docx").exists()


def test_document_capability_is_registered_and_routes_automatically(tmp_path: Path):
    assert any(cap.name == "document_editing" for cap in create_all_capabilities())
    reset_capability_registry()
    try:
        registry = CapabilityRegistry(); doc = capability(tmp_path); registry.register(doc); registry.start()
        router = CapabilityRouter(); bridge = CapabilityRegistrationBridge(registry=registry, router=router, tool_manager=ToolManager(tmp_path, file_allowlist=allowlist(tmp_path)))
        bridge.register_registered_capability(doc, patterns=[r"\bdocument\b"], keywords=["document"])
        source = tmp_path / "route.md"; source.write_text("routed", encoding="utf-8")
        result = router.route("document", capability_action="inspect", file_path=str(source))
        assert result.success is True
        assert result.capability_name == "document_editing"
    finally:
        reset_capability_registry()
