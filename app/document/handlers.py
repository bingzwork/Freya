"""Format-specific document handlers used by DocumentEditingCapability."""
from __future__ import annotations

import csv
import io
import re
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional


TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".html", ".htm"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | {".docx", ".pdf", ".xlsx", ".csv", ".pptx"}


class DocumentHandler(ABC):
    extensions: set[str] = set()

    @abstractmethod
    def inspect(self, path: Path) -> Dict[str, Any]: ...

    @abstractmethod
    def edit(self, path: Path, inputs: Dict[str, Any]) -> bytes: ...

    def export(self, path: Path, target_extension: str, inputs: Dict[str, Any]) -> bytes:
        if target_extension == path.suffix.lower():
            return self.edit(path, {**inputs, "edits": []})
        raise ValueError(f"Export from {path.suffix} to {target_extension} is not supported")

    def validate(self, content: bytes, extension: str) -> Dict[str, Any]:
        if not content:
            raise ValueError("Generated document is empty")
        return {"valid": True, "size_bytes": len(content), "extension": extension}

    @staticmethod
    def _replace(value: str, inputs: Dict[str, Any]) -> str:
        if "content" in inputs:
            return str(inputs["content"])
        result = value
        for edit in inputs.get("edits", []) or []:
            if not isinstance(edit, dict) or "old" not in edit:
                raise ValueError("Each edit must contain an 'old' value")
            old, new = str(edit["old"]), str(edit.get("new", ""))
            if bool(edit.get("regex", False)):
                result = re.sub(old, new, result, count=int(edit.get("count", 0)))
            else:
                count = int(edit.get("count", -1))
                result = result.replace(old, new, count)
        return result


class TextHandler(DocumentHandler):
    extensions = TEXT_EXTENSIONS

    def inspect(self, path: Path) -> Dict[str, Any]:
        text = path.read_text(encoding="utf-8")
        return {"format": path.suffix.lower().lstrip("."), "text": text, "line_count": len(text.splitlines()), "character_count": len(text)}

    def edit(self, path: Path, inputs: Dict[str, Any]) -> bytes:
        text = path.read_text(encoding=inputs.get("encoding", "utf-8"))
        return self._replace(text, inputs).encode(inputs.get("encoding", "utf-8"))

    def export(self, path: Path, target_extension: str, inputs: Dict[str, Any]) -> bytes:
        text = self._replace(path.read_text(encoding="utf-8"), inputs)
        if target_extension == ".html":
            return (f"<!doctype html><html><body><pre>{_escape_html(text)}</pre></body></html>").encode("utf-8")
        if target_extension in TEXT_EXTENSIONS:
            return text.encode("utf-8")
        if target_extension == ".docx":
            return DocxHandler.text_to_docx(text)
        if target_extension == ".pdf":
            return PdfHandler.text_to_pdf(text)
        raise ValueError(f"Text export to {target_extension} is not supported")


class DocxHandler(DocumentHandler):
    extensions = {".docx"}

    @staticmethod
    def _document(path: Path):
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("DOCX support requires python-docx") from exc
        return Document(str(path))

    @staticmethod
    def text_to_docx(text: str) -> bytes:
        from docx import Document
        document = Document()
        for line in text.splitlines() or [""]:
            document.add_paragraph(line)
        output = io.BytesIO(); document.save(output); return output.getvalue()

    def inspect(self, path: Path) -> Dict[str, Any]:
        document = self._document(path)
        paragraphs = [p.text for p in document.paragraphs]
        tables = [[cell.text for cell in row.cells] for table in document.tables for row in table.rows]
        return {"format": "docx", "paragraphs": paragraphs, "tables": tables, "text": "\n".join(paragraphs)}

    def edit(self, path: Path, inputs: Dict[str, Any]) -> bytes:
        document = self._document(path)
        for paragraph in document.paragraphs:
            replacement = self._replace(paragraph.text, inputs)
            if replacement != paragraph.text:
                if paragraph.runs:
                    paragraph.runs[0].text = replacement
                    for run in paragraph.runs[1:]: run.text = ""
                else:
                    paragraph.text = replacement
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    replacement = self._replace(cell.text, inputs)
                    if replacement != cell.text: cell.text = replacement
        output = io.BytesIO(); document.save(output); return output.getvalue()

    def export(self, path: Path, target_extension: str, inputs: Dict[str, Any]) -> bytes:
        if target_extension == ".docx": return self.edit(path, inputs)
        text = self.inspect(path)["text"]
        if target_extension == ".txt": return text.encode("utf-8")
        if target_extension == ".pdf": return PdfHandler.text_to_pdf(self._replace(text, inputs))
        raise ValueError(f"DOCX export to {target_extension} is not supported")


class SpreadsheetHandler(DocumentHandler):
    extensions = {".xlsx", ".csv"}

    def inspect(self, path: Path) -> Dict[str, Any]:
        if path.suffix.lower() == ".csv":
            with path.open(newline="", encoding="utf-8") as stream: rows = list(csv.reader(stream))
            return {"format": "csv", "sheets": {"Sheet1": rows}, "row_count": len(rows)}
        from openpyxl import load_workbook
        workbook = load_workbook(path, data_only=False)
        sheets = {name: [[cell.value for cell in row] for row in workbook[name].iter_rows()] for name in workbook.sheetnames}
        return {"format": "xlsx", "sheet_names": workbook.sheetnames, "sheets": sheets}

    def edit(self, path: Path, inputs: Dict[str, Any]) -> bytes:
        if path.suffix.lower() == ".csv":
            rows = list(csv.reader(path.open(newline="", encoding="utf-8")))
            for change in inputs.get("cells", []): rows[int(change["row"]) - 1][int(change["column"]) - 1] = str(change.get("value", ""))
            output = io.StringIO(); csv.writer(output, lineterminator="\n").writerows(rows); return output.getvalue().encode("utf-8")
        from openpyxl import load_workbook
        workbook = load_workbook(path, data_only=False)
        for change in inputs.get("cells", []):
            sheet = str(change.get("sheet") or workbook.sheetnames[0])
            workbook[sheet][str(change["cell"])] = change.get("value")
        for sheet in inputs.get("add_sheets", []):
            if sheet not in workbook.sheetnames: workbook.create_sheet(str(sheet))
        output = io.BytesIO(); workbook.save(output); return output.getvalue()

    def validate(self, content: bytes, extension: str) -> Dict[str, Any]:
        if extension == ".xlsx":
            from openpyxl import load_workbook
            load_workbook(io.BytesIO(content), data_only=False).close()
        else:
            content.decode("utf-8")
        return super().validate(content, extension)


class PptxHandler(DocumentHandler):
    extensions = {".pptx"}

    @staticmethod
    def _presentation(path_or_stream):
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("PPTX support requires python-pptx") from exc
        return Presentation(path_or_stream)

    def inspect(self, path: Path) -> Dict[str, Any]:
        presentation = self._presentation(str(path)); slides = []
        for slide in presentation.slides:
            slides.append([shape.text for shape in slide.shapes if hasattr(shape, "text_frame")])
        return {"format": "pptx", "slide_count": len(presentation.slides), "slides": slides}

    def edit(self, path: Path, inputs: Dict[str, Any]) -> bytes:
        presentation = self._presentation(str(path))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"): continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = self._replace(run.text, inputs)
        output = io.BytesIO(); presentation.save(output); return output.getvalue()

    def validate(self, content: bytes, extension: str) -> Dict[str, Any]:
        self._presentation(io.BytesIO(content)); return super().validate(content, extension)


class PdfHandler(DocumentHandler):
    extensions = {".pdf"}

    @staticmethod
    def text_to_pdf(text: str) -> bytes:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        output = io.BytesIO(); page = canvas.Canvas(output, pagesize=letter); width, height = letter
        y = height - 48
        for line in text.splitlines() or [""]:
            if y < 48: page.showPage(); y = height - 48
            page.drawString(48, y, line[:120]); y -= 14
        page.save(); return output.getvalue()

    def inspect(self, path: Path) -> Dict[str, Any]:
        from pypdf import PdfReader
        reader = PdfReader(str(path)); pages = [page.extract_text() or "" for page in reader.pages]
        return {"format": "pdf", "page_count": len(pages), "pages": pages, "text": "\n".join(pages)}

    def edit(self, path: Path, inputs: Dict[str, Any]) -> bytes:
        from pypdf import PdfReader, PdfWriter
        reader = PdfReader(str(path)); writer = PdfWriter()
        pages = list(reader.pages)
        if inputs.get("remove_pages"):
            remove = {int(index) for index in inputs["remove_pages"]}; pages = [page for index, page in enumerate(pages) if index not in remove]
        if inputs.get("reorder_pages"):
            pages = [pages[int(index)] for index in inputs["reorder_pages"]]
        if inputs.get("merge_paths"):
            for merge_path in inputs["merge_paths"]:
                pages.extend(PdfReader(str(merge_path)).pages)
        for page in pages: writer.add_page(page)
        output = io.BytesIO(); writer.write(output); return output.getvalue()

    def export(self, path: Path, target_extension: str, inputs: Dict[str, Any]) -> bytes:
        if target_extension == ".pdf": return self.edit(path, inputs)
        if target_extension == ".txt": return self.inspect(path)["text"].encode("utf-8")
        raise ValueError(f"PDF export to {target_extension} is not supported")

    def validate(self, content: bytes, extension: str) -> Dict[str, Any]:
        from pypdf import PdfReader
        PdfReader(io.BytesIO(content)); return super().validate(content, extension)


def _escape_html(value: str) -> str:
    return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def handler_for(path: Path) -> DocumentHandler:
    extension = path.suffix.lower()
    for handler in (TextHandler(), DocxHandler(), SpreadsheetHandler(), PptxHandler(), PdfHandler()):
        if extension in handler.extensions: return handler
    raise ValueError(f"Unsupported document format: {extension or '<none>'}")


__all__ = ["DocumentHandler", "SUPPORTED_EXTENSIONS", "handler_for"]
